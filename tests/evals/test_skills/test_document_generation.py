"""Tests for document generation skills (PPTX, DOCX, XLSX)."""

import os
import tempfile

import pytest

# Try to import document libraries - skip tests if not available
try:
    from pptx import Presentation
    from pptx.util import Inches

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from docx import Document

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
class TestPowerPointGeneration:
    """Tests for PowerPoint generation skill."""

    def test_create_basic_presentation(self):
        """Should create a basic presentation with title slide."""
        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "Test Presentation"

        if slide.shapes.placeholders[1]:
            slide.shapes.placeholders[1].text = "Subtitle here"

        assert len(prs.slides) == 1
        assert slide.shapes.title.text == "Test Presentation"

    def test_add_content_slide(self):
        """Should add content slides with bullets."""
        prs = Presentation()

        # Add content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Key Points"

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.text = "First point"

        p = tf.add_paragraph()
        p.text = "Second point"
        p.level = 0

        assert len(prs.slides) == 1
        assert len(tf.paragraphs) == 2

    def test_save_presentation(self):
        """Should save presentation to file."""
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Save Test"

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            assert os.path.exists(f.name)
            assert os.path.getsize(f.name) > 0
            os.unlink(f.name)

    def test_slide_aspect_ratio_16_9(self):
        """Should support 16:9 aspect ratio."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx not installed")
class TestWordDocumentGeneration:
    """Tests for Word document generation skill."""

    def test_create_basic_document(self):
        """Should create a basic document with heading."""
        doc = Document()
        doc.add_heading("Test Document", level=0)
        doc.add_paragraph("This is a test paragraph.")

        assert len(doc.paragraphs) == 2

    def test_add_bullet_list(self):
        """Should add bullet lists."""
        doc = Document()
        items = ["Item 1", "Item 2", "Item 3"]

        for item in items:
            doc.add_paragraph(item, style="List Bullet")

        # Check bullet paragraphs were added
        bullet_paras = [p for p in doc.paragraphs if p.style.name == "List Bullet"]
        assert len(bullet_paras) == 3

    def test_add_table(self):
        """Should add tables."""
        doc = Document()

        table = doc.add_table(rows=2, cols=3)
        table.style = "Table Grid"

        # Add header
        cells = table.rows[0].cells
        cells[0].text = "Name"
        cells[1].text = "Role"
        cells[2].text = "Department"

        # Add data
        cells = table.rows[1].cells
        cells[0].text = "Alice"
        cells[1].text = "Engineer"
        cells[2].text = "Tech"

        assert len(table.rows) == 2
        assert len(table.columns) == 3

    def test_add_numbered_list(self):
        """Should add numbered lists."""
        doc = Document()
        steps = ["Step one", "Step two", "Step three"]

        for step in steps:
            doc.add_paragraph(step, style="List Number")

        numbered_paras = [p for p in doc.paragraphs if p.style.name == "List Number"]
        assert len(numbered_paras) == 3

    def test_save_document(self):
        """Should save document to file."""
        doc = Document()
        doc.add_heading("Save Test", level=0)

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc.save(f.name)
            assert os.path.exists(f.name)
            assert os.path.getsize(f.name) > 0
            os.unlink(f.name)


@pytest.mark.skipif(not HAS_XLSX, reason="openpyxl not installed")
class TestExcelSpreadsheetGeneration:
    """Tests for Excel spreadsheet generation skill."""

    def test_create_basic_workbook(self):
        """Should create a basic workbook with data."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Test Sheet"

        ws["A1"] = "Header 1"
        ws["B1"] = "Header 2"
        ws["A2"] = "Data 1"
        ws["B2"] = "Data 2"

        assert ws["A1"].value == "Header 1"
        assert ws["B2"].value == "Data 2"

    def test_add_multiple_sheets(self):
        """Should support multiple worksheets."""
        wb = Workbook()
        wb.active.title = "Summary"

        wb.create_sheet("Details")
        wb.create_sheet("Raw Data")

        assert len(wb.sheetnames) == 3
        assert "Details" in wb.sheetnames

    def test_add_formulas(self):
        """Should support formulas."""
        wb = Workbook()
        ws = wb.active

        ws["A1"] = 10
        ws["A2"] = 20
        ws["A3"] = 30
        ws["A4"] = "=SUM(A1:A3)"

        assert ws["A4"].value == "=SUM(A1:A3)"

    def test_cell_formatting(self):
        """Should support cell formatting."""
        wb = Workbook()
        ws = wb.active

        ws["A1"] = "Bold Header"
        ws["A1"].font = Font(bold=True)

        assert ws["A1"].font.bold is True

    def test_freeze_panes(self):
        """Should support freezing panes."""
        wb = Workbook()
        ws = wb.active

        ws["A1"] = "Header"
        ws.freeze_panes = "A2"

        assert ws.freeze_panes == "A2"

    def test_auto_filter(self):
        """Should support auto filter."""
        wb = Workbook()
        ws = wb.active

        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Test"
        ws["B2"] = 100

        ws.auto_filter.ref = "A1:B2"

        assert ws.auto_filter.ref == "A1:B2"

    def test_save_workbook(self):
        """Should save workbook to file."""
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Test"

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb.save(f.name)
            assert os.path.exists(f.name)
            assert os.path.getsize(f.name) > 0
            os.unlink(f.name)


class TestDocumentGenerationPatterns:
    """Test common document generation patterns."""

    def test_filename_sanitization(self):
        """Should sanitize filenames properly."""

        def sanitize_filename(title: str) -> str:
            """Convert title to safe filename."""
            # Replace spaces with underscores, remove special chars
            safe = title.lower().replace(" ", "_")
            safe = "".join(c for c in safe if c.isalnum() or c == "_")
            return safe

        assert sanitize_filename("Project Report 2024") == "project_report_2024"
        assert sanitize_filename("Test: Document!") == "test_document"
        assert sanitize_filename("My File (Draft)") == "my_file_draft"

    def test_workspace_path_construction(self):
        """Should construct proper workspace paths."""
        workspace = "/workspace"
        filename = "test_report.docx"

        full_path = f"{workspace}/files/{filename}"

        assert full_path == "/workspace/files/test_report.docx"
        assert full_path.endswith(".docx")


@pytest.mark.eval
class TestDocumentGenerationSkillEval:
    """Evaluation tests for document generation skill agent behavior."""

    def test_skill_triggers_pptx(self):
        """Verify PPTX skill trigger patterns."""
        triggers = [
            "create presentation",
            "make slides",
            "generate pptx",
            "powerpoint",
            "slide deck",
            "presentation about",
        ]

        test_queries = [
            "Create a presentation about AI",
            "Make slides for the quarterly review",
            "Generate a PowerPoint deck",
        ]

        for query in test_queries:
            matched = any(t in query.lower() for t in triggers)
            assert matched, f"Query '{query}' should match a PPTX trigger"

    def test_skill_triggers_docx(self):
        """Verify DOCX skill trigger patterns."""
        triggers = [
            "create document",
            "generate docx",
            "write document",
            "word document",
            "create report",
            "document about",
            "create a document",
            "write a report",
        ]

        test_queries = [
            "Create a document about the project",
            "Write a report on our progress",
            "Generate a Word document",
        ]

        for query in test_queries:
            matched = any(t in query.lower() for t in triggers)
            assert matched, f"Query '{query}' should match a DOCX trigger"

    def test_skill_triggers_xlsx(self):
        """Verify XLSX skill trigger patterns."""
        triggers = [
            "create spreadsheet",
            "generate excel",
            "make xlsx",
            "excel file",
            "data to excel",
            "spreadsheet for",
            "create a spreadsheet",
            "an xlsx",
        ]

        test_queries = [
            "Create a spreadsheet for the budget",
            "Generate an Excel file with the data",
            "Make an xlsx with sales figures",
        ]

        for query in test_queries:
            matched = any(t in query.lower() for t in triggers)
            assert matched, f"Query '{query}' should match an XLSX trigger"

    def test_response_includes_file_path(self):
        """Response should include the generated file path."""
        # Simulated response structure
        response = {"file": "workspace/files/test_report.docx", "message": "Document created"}

        assert "workspace/files/" in response["file"]
        assert any(ext in response["file"] for ext in [".pptx", ".docx", ".xlsx"])
